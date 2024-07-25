import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import os
import json
import pandas as pd
import numpy as np
from numpy.linalg import norm
import hashlib
import openai
from dotenv import load_dotenv
load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")
DATA_DIR = 'C:\\Users\\user\\ahhshit\\temp'
embeddings = []
model_name = "all-MiniLM-L6-v2"  # Replace with your desired model name
model = SentenceTransformer(model_name)

def fetch_user_dict(UserID):
    user_data_file = os.path.join(DATA_DIR, f'{UserID}.json')
    if os.path.exists(user_data_file):
        with open(user_data_file, 'r') as f:
            return json.load(f)
    else:
        # Create a new user file with the name UserID.json
        with open(user_data_file, 'w') as f:
            json.dump({}, f)
        return {}
   

def save_user_dict(UserID, data):
    user_data_file = os.path.join(DATA_DIR, f'{UserID}.json')
    with open(user_data_file, 'w') as f:
        json.dump(data, f)

def post_embeddings(UserID, BookId, pageno, pagetext, embeddings):
    user_dict = fetch_user_dict(UserID)
    if BookId not in user_dict:
        user_dict[BookId] = {}
    if pageno not in user_dict[BookId]:
        user_dict[BookId][pageno] = {}
    user_dict[BookId][pageno][pagetext] = embeddings
    print("Embeddings added successfully\n")
    save_user_dict(UserID, user_dict)

def retrieve_topk(query_embeddings, UserID, top_k):
    user_dict = fetch_user_dict(UserID)
    similarities = []
    query_embeddings = np.array(query_embeddings, dtype=float)
    for BookId, pageno in user_dict.items():
        for pagenum, page_content in pageno.items():
            for text, embeddings in page_content.items():
                embeddings = np.array(embeddings, dtype=float)
                if embeddings.ndim > 1:
                    embeddings = embeddings.squeeze()
                similarity_score = np.dot(query_embeddings, embeddings) / (norm(query_embeddings) * norm(embeddings))
                similarities.append((BookId, pagenum, text, embeddings, similarity_score))
    similarities.sort(key=lambda x: x[4], reverse=True)
    top_k_results = similarities[:top_k]
    return top_k_results

def get_user_books(UserID):
    user_dict = fetch_user_dict(UserID)
    return [bookid for bookid in user_dict]

def del_user_book(UserID, BookID):
    user_dict = fetch_user_dict(UserID)
    if BookID in user_dict:
        del user_dict[BookID]
        save_user_dict(UserID, user_dict)
        print(f"Book '{BookID}' removed from library.")
    else:
        print(f"Book '{BookID}' does not exist in the library.")

def get_embeddings(text):
    if not isinstance(text, (list, str)):
        raise TypeError("Input text must be a string or a list of strings.")
    if isinstance(text, str):
        text = [text]
    embeddings = model.encode(text)
    return embeddings

def read_content(file_path, file_name, UserID):
    if file_path.endswith(".pdf"):
        with open(file_path, 'rb') as pdf_file:
            reader = PdfReader(pdf_file)
            num_pages = len(reader.pages)
            for page_num in range(num_pages):
                page = reader.pages[page_num]
                text = page.extract_text()
                if text:
                    page_embeddings = get_embeddings(text)
                    embeddings = page_embeddings.tolist()
                    post_embeddings(UserID, file_name, page_num, text, embeddings)
    elif file_path.endswith(".docx"):
        document = Document(file_path)
        for page_num, page in enumerate(document.paragraphs):
            text = page.text.strip()
            if text:
                page_embeddings = get_embeddings(text)
                embeddings = page_embeddings.tolist()
                post_embeddings(UserID, file_name, page_num, text, embeddings)
    elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text
            if text:
                page_embeddings = get_embeddings(text)
                embeddings = page_embeddings.tolist()
                post_embeddings(UserID, file_name, slide_num, text, embeddings)
    else:
        print(f"Unsupported file format: {file_path}")

def get_embeddings_by_book(UserID, BookID):
    user_dict = fetch_user_dict(UserID)
    return user_dict.get(BookID, {})

def query_embedding(query):
    return model.encode(query)

def call_openai_api(prompt):
    openai.api_key = openai_api_key
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=3000,
        temperature=0.9,
        top_p=1.0,
        frequency_penalty=0.0,
        presence_penalty=0.0
    )
    return response.choices[0].message['content']

def relevant_results(results):
    relevant_content = ""
    for bookid, pageno, text, embeddings, similarity_score in results:
        st.write("bookname: ",bookid)
        st.write("page number:",pageno)
        st.write("score:",similarity_score)
        relevant_content += text + "\n"
    print("result:",relevant_results)
    return relevant_content

def generate_user_id(email):
    return hashlib.md5(email.encode()).hexdigest()

def authenticate(email, password):
    auth_data_file = os.path.join(DATA_DIR, 'auth_data.json')
    if os.path.exists(auth_data_file):
        with open(auth_data_file, 'r') as f:
            auth_data = json.load(f)
        if email in auth_data and auth_data[email]['password'] == hashlib.md5(password.encode()).hexdigest():
            return auth_data[email]['user_id']
    return None

def register_user(email, password):
    auth_data_file = os.path.join(DATA_DIR, 'auth_data.json')
    if os.path.exists(auth_data_file):
        with open(auth_data_file, 'r') as f:
            auth_data = json.load(f)
    else:
        auth_data = {}
    if email in auth_data:
        return None
    user_id = generate_user_id(email)
    auth_data[email] = {
        'user_id': user_id,
        'password': hashlib.md5(password.encode()).hexdigest()
    }
    with open(auth_data_file, 'w') as f:
        json.dump(auth_data, f)
    return user_id

def main():
    st.title("Document Processing and Query System")

    if "user_id" not in st.session_state:
        st.sidebar.header("Login/Register")
        email = st.sidebar.text_input("Email")
        password = st.sidebar.text_input("Password", type="password")

        if st.sidebar.button("Login"):
            user_id = authenticate(email, password)
            if user_id:
                st.session_state.user_id = user_id
                st.success("Logged in successfully!")
            else:
                st.error("Invalid email or password")

        if st.sidebar.button("Register"):
            user_id = register_user(email, password)
            if user_id:
                st.session_state.user_id = user_id
                st.success("Registered successfully! Please log in.")
            else:
                st.error("Email already registered")

    if "user_id" in st.session_state:
        UserID = st.session_state.user_id
        st.sidebar.header("Upload your documents")
        uploaded_files = st.sidebar.file_uploader("Choose files", accept_multiple_files=True, type=['pdf', 'docx', 'pptx'])

        if uploaded_files:
            
            
            for uploaded_file in uploaded_files:
    
                with st.spinner(f'Processing {uploaded_file.name}...'):
                    temp_dir = "temp"
                    os.makedirs(temp_dir, exist_ok=True)
                    file_name = uploaded_file.name
                    file_path = os.path.join(temp_dir, uploaded_file.name)
                    with open(file_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    read_content(file_path, file_name, UserID)
                    st.write("Document processed successfully!")

        st.sidebar.subheader("Select a document to chat with")
        document_names = get_user_books(UserID)
        selected_doc = st.sidebar.selectbox("Available documents", document_names)

        if selected_doc:
            st.write(f"You selected: {selected_doc}")
            query = st.text_input(f"Enter your query for {selected_doc}:")
            if query:
                with st.spinner('Processing query...'):
                    query_embeddings = query_embedding(query)
                    results = retrieve_topk(query_embeddings, UserID, top_k=3)
                    relevant_content = relevant_results(results)
                    prompt = f"Based on the following context, answer the question: {relevant_content}\n\nQuestion: {query}\nAnswer:"
                    result = call_openai_api(prompt)
                    st.write(result)

if __name__ == '__main__':
    main()
